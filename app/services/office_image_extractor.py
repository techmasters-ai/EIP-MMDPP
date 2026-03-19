"""Extract embedded images from Office documents (DOCX, PPTX).

Docling's SimplePipeline detects PictureItem positions but doesn't extract
the actual image bytes from Office formats. This module fills that gap by
using python-docx and python-pptx to extract images, then injecting them
into the Docling JSON so derive_picture_descriptions can process them.
"""

import base64
import io
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def extract_docx_images(file_bytes: bytes) -> list[dict]:
    """Extract all images from a DOCX file.

    Returns list of dicts: [{"index": 0, "b64": "...", "content_type": "image/png"}, ...]
    """
    from docx import Document

    images = []
    try:
        doc = Document(io.BytesIO(file_bytes))
        for idx, rel in enumerate(doc.part.rels.values()):
            if "image" in rel.reltype:
                try:
                    image_part = rel.target_part
                    image_bytes = image_part.blob
                    content_type = image_part.content_type or "image/png"
                    b64 = base64.b64encode(image_bytes).decode("ascii")
                    images.append({
                        "index": idx,
                        "b64": b64,
                        "content_type": content_type,
                    })
                except Exception as e:
                    logger.debug("Failed to extract DOCX image %d: %s", idx, e)
    except Exception as e:
        logger.warning("Failed to open DOCX for image extraction: %s", e)

    logger.info("Extracted %d images from DOCX", len(images))
    return images


def extract_pptx_images(file_bytes: bytes) -> list[dict]:
    """Extract all images from a PPTX file.

    Returns list of dicts: [{"index": 0, "b64": "...", "content_type": "image/png", "slide": 1}, ...]
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    images = []
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        img_idx = 0
        for slide_num, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        content_type = image.content_type or "image/png"
                        b64 = base64.b64encode(image_bytes).decode("ascii")
                        images.append({
                            "index": img_idx,
                            "b64": b64,
                            "content_type": content_type,
                            "slide": slide_num,
                        })
                        img_idx += 1
                    except Exception as e:
                        logger.debug("Failed to extract PPTX image on slide %d: %s", slide_num, e)
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    # Check grouped shapes for pictures
                    for child in shape.shapes:
                        if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            try:
                                image = child.image
                                image_bytes = image.blob
                                content_type = image.content_type or "image/png"
                                b64 = base64.b64encode(image_bytes).decode("ascii")
                                images.append({
                                    "index": img_idx,
                                    "b64": b64,
                                    "content_type": content_type,
                                    "slide": slide_num,
                                })
                                img_idx += 1
                            except Exception as e:
                                logger.debug("Failed to extract grouped PPTX image: %s", e)
    except Exception as e:
        logger.warning("Failed to open PPTX for image extraction: %s", e)

    logger.info("Extracted %d images from PPTX", len(images))
    return images


def inject_images_into_docling_json(
    docling_json: dict,
    images: list[dict],
) -> int:
    """Inject extracted image data into Docling JSON PictureItem entries.

    Matches extracted images to PictureItem entries that are missing image data.
    Modifies docling_json in place.

    Returns number of images injected.
    """
    # Find picture entries missing image data in the "pictures" collection
    pictures = docling_json.get("pictures", [])
    if not isinstance(pictures, list):
        return 0

    empty_pics = []
    for pic in pictures:
        if not isinstance(pic, dict):
            continue
        image_ref = pic.get("image", {})
        if not isinstance(image_ref, dict):
            image_ref = {}
        uri = image_ref.get("uri", "")
        # Picture has no image data
        if not uri or not uri.startswith("data:"):
            empty_pics.append(pic)

    if not empty_pics or not images:
        return 0

    # Match by order: first empty picture gets first extracted image, etc.
    injected = 0
    for pic, img in zip(empty_pics, images):
        ct = img["content_type"]
        b64 = img["b64"]
        data_uri = f"data:{ct};base64,{b64}"

        if "image" not in pic:
            pic["image"] = {}
        pic["image"]["uri"] = data_uri

        # Set image dimensions if possible
        try:
            from PIL import Image
            pil_img = Image.open(io.BytesIO(base64.b64decode(b64)))
            pic["image"]["size"] = {"width": pil_img.width, "height": pil_img.height}
        except Exception:
            pass

        injected += 1

    logger.info("Injected %d images into Docling JSON (%d empty slots, %d images available)",
                injected, len(empty_pics), len(images))
    return injected
